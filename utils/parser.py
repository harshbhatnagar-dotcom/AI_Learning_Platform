import fitz
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
import io
from dotenv import load_dotenv
from mistralai.client import Mistral

load_dotenv(override=True)

import os


client = Mistral(api_key=os.environ["MISTRAL_API_KEY"])


def mistral_ocr(uploaded_file):
    """
    Perform OCR on a PDF using Mistral OCR.
    """

    uploaded_file.seek(0)

    uploaded = client.files.upload(
        file={
            "file_name": uploaded_file.name,
            "content": uploaded_file.read(),
        },
        purpose="ocr",
    )

    response = client.ocr.process(
        model="mistral-ocr-latest",
        document={
            "type": "file",
            "file_id": uploaded.id,
        },
        include_image_base64=False,
    )

    pages = []

    for page in response.pages:
        pages.append(page.markdown)

    # Optional: delete uploaded file from Mistral
    client.files.delete(uploaded.id)

    return "\n\n".join(pages)

def tesseract_ocr(pdf_bytes):
    """
    Perform OCR on a PDF using Tesseract.
    """
    doc = fitz.open(
        stream=pdf_bytes,
        filetype="pdf"
    )

    ocr_text = []

    for page in doc:
        pix = page.get_pixmap(dpi=300)
        img = Image.open(
            io.BytesIO(
                pix.tobytes("png")
            )
        )
        extracted = pytesseract.image_to_string(
            img,
            lang="eng",
            config="--oem 3 --psm 6"
        )
        ocr_text.append(extracted)

    return "\n".join(ocr_text)


def extract_pdf(uploaded_file):
    """
    Extract text from a PDF.
    Falls back to OCR if no text is found.
    """

    uploaded_file.seek(0)

    pdf_bytes = uploaded_file.read()

    doc = fitz.open(
        stream=pdf_bytes,
        filetype="pdf"
    )

    text = ""

    for page in doc:
        text += page.get_text()

    text = text.strip()

    if text:
        return text
    
    uploaded_file.seek(0)
    return mistral_ocr(uploaded_file)

    


def extract_docx(uploaded_file):
    """
    Extract text from DOCX.
    """

    uploaded_file.seek(0)

    doc = Document(uploaded_file)

    text = ""

    for para in doc.paragraphs:
        text += para.text + "\n"

    return text


def extract_pptx(uploaded_file):
    """
    Extract text from PPTX.
    """

    uploaded_file.seek(0)

    prs = Presentation(uploaded_file)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def extract_text(uploaded_file):
    """
    Automatically detect the uploaded file type
    and extract text.
    """

    filename = uploaded_file.name.lower()

    if filename.endswith(".pdf"):
        return extract_pdf(uploaded_file)

    elif filename.endswith(".docx"):
        return extract_docx(uploaded_file)

    elif filename.endswith(".pptx"):
        return extract_pptx(uploaded_file)

    else:
        raise ValueError("Unsupported file format.")